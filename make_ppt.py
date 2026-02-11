from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# === CREATE PRESENTATION ===
prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    # uniform font style
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(22)
                run.font.color.rgb = RGBColor(0, 51, 102)
    return slide

# === SLIDE 1: Cover ===
cover = prs.slides.add_slide(prs.slide_layouts[6])
txBox = cover.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(4))
tf = txBox.text_frame
tf.text = "Medi-Caps University, Indore\n\nDepartment of Computer Science and Engineering\n\n"
p = tf.add_paragraph()
p.text = "Project Report On\nDROWSINESS DETECTION SYSTEM USING COMPUTER VISION"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p = tf.add_paragraph()
p.text = "\nSubmitted By:\nJignesh Savla (EN22CS303025)\nLucky Yadav (EN22CS303030)\n\nGuides:\nMandakini Ingle, Rahul Choudhary"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 51, 102)

# === SLIDE 2-10 ===
add_slide("Introduction",
           "• Detect driver drowsiness using real-time video.\n"
           "• Prevent accidents by alerting the driver.\n"
           "• Uses computer vision and machine learning.")

add_slide("Objectives",
           "• Detect closed eyes or yawning.\n"
           "• Generate alert sound when fatigue is detected.\n"
           "• Improve road safety through automation.")

add_slide("Existing System",
           "• Manual observation or simple alarm systems.\n"
           "• Lack of automation.\n"
           "• High risk due to delayed response.")

add_slide("Proposed System",
           "• Detects drowsiness through webcam.\n"
           "• Monitors eyes, yawning, and face position.\n"
           "• Uses OpenCV, Dlib, and Python.\n"
           "• Alerts in real-time.")

add_slide("System Architecture",
           "• Input: Camera feed.\n"
           "• Processing: Face and eye detection using ML.\n"
           "• Output: Alarm if driver is drowsy.\n"
           "• Database: Logs events for review.")

add_slide("Database Design",
           "• Normalized structure.\n"
           "• Handles data redundancy.\n"
           "• ER diagram connects all entities clearly.")

add_slide("GUI Design",
           "• Start/Stop buttons for monitoring.\n"
           "• Real-time display.\n"
           "• User-friendly and minimalistic design.")

add_slide("Results",
           "• Accurate eye closure detection.\n"
           "• Fast alert generation.\n"
           "• Improved accuracy and usability.")

add_slide("Conclusion",
           "• Drowsiness detection helps prevent accidents.\n"
           "• Enhances driver safety.\n"
           "• Can be extended with IoT and mobile alerts.")

add_slide("Thank You",
           "Medi-Caps University, Indore\nDepartment of Computer Science and Engineering")

# === SAVE FILE ===
prs.save("Drowsiness_Detection_System_Presentation.pptx")
print("✅ Presentation Created: Drowsiness_Detection_System_Presentation.pptx")